from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    # 1. Initialize the presentation
    prs = Presentation()

    # --- SLIDE 1: Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "The Daily Canvas OS"
    subtitle1.text = "A Modular, AI-Powered Productivity Dashboard\nDeveloped by Eldho"

    # --- SLIDE 2: System Architecture ---
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    slide2.shapes.title.text = "System Architecture & Tech Stack"
    body_shape2 = slide2.shapes.placeholders[1]
    tf2 = body_shape2.text_frame
    tf2.text = "Built entirely with Vanilla JavaScript (ES6+), HTML5, and CSS3."
    p = tf2.add_paragraph()
    p.text = "API Integrations: Google Gemini 2.5 (LLM), Open-Meteo, CoinGecko, and Wikimedia."
    p.level = 1
    p2 = tf2.add_paragraph()
    p2.text = "DOM Physics Engine: Custom JS mathematics for 20px 'Snap-to-Grid' draggable widgets."
    p2.level = 1

    # --- SLIDE 3: UX & Visual Innovations ---
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    slide3.shapes.title.text = "UI/UX & Visual Engineering"
    tf3 = slide3.shapes.placeholders[1].text_frame
    tf3.text = "The 'Chameleon' Visual Engine"
    p = tf3.add_paragraph()
    p.text = "Utilizes HTML5 Canvas and ColorThief.js to dynamically extract RGB values from background art."
    p.level = 1
    p2 = tf3.add_paragraph()
    p2.text = "Zen Mode: Hardware-toggled (Esc key) UI decluttering for deep focus."
    p2.level = 0
    p3 = tf3.add_paragraph()
    p3.text = "Glassmorphism UI: CSS backdrop-filters for a modern, sleek aesthetic."
    p3.level = 1

    # --- SLIDE 4: Academic Power Features ---
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    slide4.shapes.title.text = "Student-Centric Power Features"
    tf4 = slide4.shapes.placeholders[1].text_frame
    tf4.text = "Smart Command Router"
    p = tf4.add_paragraph()
    p.text = "Intercepts '/ai' for Gemini queries and '/t' for task logging directly from the search bar."
    p.level = 1
    p2 = tf4.add_paragraph()
    p2.text = "Active Recall Flashcards: State-managed study tool with 3D CSS flip animations."
    p2.level = 0
    p3 = tf4.add_paragraph()
    p3.text = "Floating Notes: Draggable text editors with one-click .txt file exporting."
    p3.level = 0

    # --- SLIDE 5: Data Persistence & Analytics ---
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    slide5.shapes.title.text = "Data Persistence & Portability"
    tf5 = slide5.shapes.placeholders[1].text_frame
    tf5.text = "Zero-Database Architecture"
    p = tf5.add_paragraph()
    p.text = "Utilizes browser localStorage to serialize arrays and objects (Tasks, Settings, Layout)."
    p.level = 1
    p2 = tf5.add_paragraph()
    p2.text = "Consistency Heatmap: Visual data matrix rendering daily activity."
    p2.level = 0
    p3 = tf5.add_paragraph()
    p3.text = "System Export: Converts all local state memory into a downloadable JSON backup via the Blob API."
    p3.level = 0

    # Save the file
    prs.save('Daily_Canvas_OS_Presentation.pptx')
    print("✅ Success! 'Daily_Canvas_OS_Presentation.pptx' has been generated.")

if __name__ == "__main__":
    create_presentation()